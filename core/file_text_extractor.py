"""
File Text Extractor - Extract text content from various file formats.

Replicates Claude's file processing logic for accurate token counting.
Supports office formats, plaintext, and archives.
"""

import logging
from pathlib import Path
from typing import Optional
import zipfile
import io

logger = logging.getLogger(__name__)

# Optional imports - graceful degradation if not available
try:
    import openpyxl
except ImportError:
    openpyxl = None
    logger.warning("openpyxl not available - Excel extraction disabled")

try:
    from docx import Document
except ImportError:
    Document = None
    logger.warning("python-docx not available - Word extraction disabled")

try:
    from pptx import Presentation
except ImportError:
    Presentation = None
    logger.warning("python-pptx not available - PowerPoint extraction disabled")


class FileTextExtractor:
    """Extract text content from files matching Claude's processing logic."""

    # Plaintext file extensions
    PLAINTEXT_EXTENSIONS = {
        '.txt', '.md', '.csv', '.json', '.xml', '.yaml', '.yml',
        '.py', '.js', '.ts', '.tsx', '.jsx', '.java', '.cpp', '.c', '.h',
        '.rs', '.go', '.rb', '.php', '.sh', '.bash', '.css', '.html', '.htm',
        '.sql', '.r', '.m', '.swift', '.kt', '.scala', '.pl', '.lua'
    }

    @staticmethod
    def extract(local_path: str, file_bytes: bytes) -> Optional[str]:
        """
        Extract text content from file.

        Args:
            local_path: Path to file (for extension detection)
            file_bytes: File content as bytes

        Returns:
            Extracted text content, or None if extraction fails
        """
        ext = Path(local_path).suffix.lower()

        try:
            if ext in {'.xlsx', '.xlsm'}:
                return FileTextExtractor._extract_excel(file_bytes)
            elif ext == '.docx':
                return FileTextExtractor._extract_word(file_bytes)
            elif ext == '.pptx':
                return FileTextExtractor._extract_powerpoint(file_bytes)
            elif ext in FileTextExtractor.PLAINTEXT_EXTENSIONS:
                return FileTextExtractor._extract_plaintext(file_bytes)
            elif ext == '.zip':
                return FileTextExtractor._extract_zip_listing(file_bytes)
            else:
                logger.debug(f"No extraction method for {ext}")
                return None

        except Exception as e:
            logger.error(f"Failed to extract text from {local_path}: {e}")
            return None

    @staticmethod
    def _extract_excel(file_bytes: bytes) -> Optional[str]:
        """
        Extract text from Excel file using openpyxl.

        Matches Claude's extraction: sheet names + cell values.
        Uses data_only=True to get calculated values, not formulas.
        """
        if openpyxl is None:
            logger.warning("openpyxl not available for Excel extraction")
            return None

        try:
            # Load from bytes with data_only=True (calculated values)
            wb = openpyxl.load_workbook(
                filename=io.BytesIO(file_bytes),
                data_only=True
            )

            parts = []
            for sheet in wb.worksheets:
                parts.append(f"Sheet: {sheet.title}")

                # Extract all non-empty rows
                for row in sheet.iter_rows(values_only=True):
                    row_values = [str(cell) for cell in row if cell is not None]
                    if row_values:
                        parts.append(" | ".join(row_values))

            return "\n".join(parts)

        except Exception as e:
            logger.error(f"Excel extraction failed: {e}")
            return None

    @staticmethod
    def _extract_word(file_bytes: bytes) -> Optional[str]:
        """
        Extract text from Word document using python-docx.

        Matches Claude's extraction: paragraphs + tables.
        """
        if Document is None:
            logger.warning("python-docx not available for Word extraction")
            return None

        try:
            doc = Document(io.BytesIO(file_bytes))
            parts = []

            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text)

            # Extract tables (critical for accurate token counts)
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    if any(row_text):
                        parts.append(" | ".join(row_text))

            return "\n".join(parts)

        except Exception as e:
            logger.error(f"Word extraction failed: {e}")
            return None

    @staticmethod
    def _extract_powerpoint(file_bytes: bytes) -> Optional[str]:
        """
        Extract text from PowerPoint using python-pptx.

        Matches Claude's extraction: slides + text shapes + speaker notes.
        """
        if Presentation is None:
            logger.warning("python-pptx not available for PowerPoint extraction")
            return None

        try:
            prs = Presentation(io.BytesIO(file_bytes))
            parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                parts.append(f"\n=== Slide {slide_num} ===")

                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text)

                # Extract speaker notes (IMPORTANT: Claude reads these!)
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes.strip():
                        parts.append(f"[Speaker Notes: {notes}]")

            return "\n".join(parts)

        except Exception as e:
            logger.error(f"PowerPoint extraction failed: {e}")
            return None

    @staticmethod
    def _extract_plaintext(file_bytes: bytes) -> Optional[str]:
        """
        Extract text from plaintext files via UTF-8 decoding.

        Handles CSV, JSON, XML, YAML, code files, etc.
        """
        try:
            # Decode with error replacement for invalid UTF-8
            return file_bytes.decode('utf-8', errors='replace')

        except Exception as e:
            logger.error(f"Plaintext extraction failed: {e}")
            return None

    @staticmethod
    def _extract_zip_listing(file_bytes: bytes) -> Optional[str]:
        """
        Extract file listing from ZIP archive.

        For now, only lists contents (not extraction).
        Full extraction may be added later.
        """
        try:
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
                file_list = zf.namelist()
                return "\n".join(file_list)

        except Exception as e:
            logger.error(f"ZIP listing extraction failed: {e}")
            return None
